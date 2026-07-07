"""
HRMS Training PPT Generator
Takes screenshots of all modules and builds a PowerPoint presentation.
"""
import os
import sys
import time
import subprocess
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from playwright.sync_api import sync_playwright

BASE_URL = 'http://127.0.0.1:5000'
OUT_DIR = 'ppt_screenshots'
PPT_OUT = 'HRMS_Department_Wise.pptx'

os.makedirs(OUT_DIR, exist_ok=True)

ADMIN_EMP = 'EMP001'
ADMIN_PASS = 'pass123'

PAGES = [
    ('/', 'login', 'Login Page', 'Employee login with credentials table showing all users and their passwords (pass123).'),
    ('/dashboard', 'admin_dashboard', 'Admin Dashboard', 'Real-time break monitoring dashboard with employee metrics, online status, and active breaks.'),
    ('/admin/users', 'admin_users', 'Employee Master', 'User management with department, designation, phone, DOB, DOJ, address, emergency contacts.'),
    ('/admin/org-chart', 'org_chart', 'Organisation Chart', 'Hierarchical org tree built from manager_emp_id relationships.'),
    ('/admin/holidays', 'holidays', 'Holiday Calendar', 'CRUD for company holidays with National/Optional types.'),
    ('/admin/import-users', 'import_users', 'CSV Import Wizard', 'Drag-and-drop employee import with sample CSV download.'),
    ('/admin/assets', 'assets', 'Asset Management', 'Issue, track, and return company assets (laptops, monitors, phones).'),
    ('/admin/jobs', 'jobs', 'Job Postings', 'Create and manage job openings with department, location, and status.'),
    ('/admin/candidates', 'candidates', 'Candidate Management', 'Track candidates through hiring pipeline: Applied Screened Interviewed Offered Hired.'),
    ('/admin/payroll', 'payroll', 'Payroll Runs', 'Monthly payroll generation with auto-calculated PF, ESI, PT deductions.'),
    ('/admin/salary-structures', 'salary', 'Salary Structures', 'Configure employee salary components: Basic, HRA, Allowances, Deductions.'),
    ('/admin/goals', 'goals', 'Performance Goals (Admin)', 'Set and rate employee goals/KRAs with weight and 1-5 star ratings.'),
    ('/admin/reviews', 'reviews', 'Performance Reviews (Admin)', 'Create review cycles, assign reviewers, track submitted ratings.'),
    ('/admin/expenses', 'expenses_admin', 'Expense Management (Admin)', 'Approve/reject employee expense claims with Paid workflow.'),
    ('/admin/tickets', 'tickets_admin', 'Help Desk (Admin)', 'Manage support tickets with comments, priority, status tracking, and assignment.'),
    ('/admin/documents', 'documents_admin', 'Document Manager (Admin)', 'Upload, download, and manage employee documents with categories.'),
    ('/admin/analytics', 'analytics', 'Advanced Analytics', 'Dashboard with 4 Chart.js charts: headcount, leave trends, expenses, performance + attrition risk.'),
    ('/leaves', 'leaves', 'Leave Management (Employee)', 'Apply for leave, view balance, track approval status with real-time notifications.'),
    ('/regularization', 'regularization', 'Regularization (Employee)', 'Submit attendance regularization requests for missing or incorrect logs.'),
    ('/onboarding', 'onboarding', 'Onboarding (Employee)', 'View assigned onboarding tasks with status tracking.'),
    ('/goals', 'my_goals', 'My Goals (Employee)', 'Set personal goals and view managerial ratings.'),
    ('/expenses', 'expenses_emp', 'My Expenses (Employee)', 'Submit expense claims with category and amount.'),
    ('/tickets', 'tickets_emp', 'My Tickets (Employee)', 'Create and track support tickets.'),
    ('/documents', 'documents_emp', 'My Documents (Employee)', 'Upload and manage personal documents.'),
    ('/profile', 'profile', 'Employee Profile', 'View/edit personal details, emergency contacts, and dependents.'),
    ('/admin/audit', 'audit', 'Audit Log', 'Track all user actions with pagination (LOGIN, LEAVE_APPLY, USER_CREATE, etc.).'),
    ('/admin/reports', 'reports', 'Reports & Export', 'Date-range reports with CSV/Excel export and employee efficiency metrics.'),
]

TABS = [
    ('/dashboard', 'admin_dashboard_productivity', 'Admin Dashboard Agent Productivity',
     'Live monitoring, dispose breaks, break summary tables with 5s auto-refresh.',
     '#productivity-tab', 'admin'),
    ('/dashboard', 'user_dashboard_modules', 'User Dashboard Modules',
     'Module cards grid: Leaves, Regularization, Goals, Expenses, Tickets, Documents, Profile.',
     '#modules-tab', 'employee'),
    ('/dashboard', 'user_dashboard_breaks', 'User Dashboard Breaks',
     'Start/end breaks, active break timer, today break history.',
     '#breaktab', 'employee'),
    ('/dashboard', 'user_dashboard_activity', 'User Dashboard My Activity',
     'Login sessions, upcoming holidays, notifications panel.',
     '#activitytab', 'employee'),
]

def capture_all():
    print("Starting Flask app...")
    proc = subprocess.Popen(
        [sys.executable, 'app.py'],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        cwd=os.path.dirname(os.path.abspath(__file__))
    )
    time.sleep(4)

    screenshots = []

    def take(page, route, fname, title, desc):
        url = f'{BASE_URL}{route}'
        print(f"  Capturing {url}...")
        try:
            page.goto(url, timeout=10000, wait_until='networkidle')
            time.sleep(1.5)
            path = os.path.join(OUT_DIR, f'{fname}.png')
            page.screenshot(path=path, full_page=True)
            screenshots.append((path, fname, title, desc))
        except Exception as e:
            print(f"  WARN: Failed {url}: {e}")
            screenshots.append((None, fname, title, desc))

    def login(page, emp_id, password):
        page.goto(f'{BASE_URL}/')
        time.sleep(0.5)
        page.fill('#empId', emp_id)
        page.fill('#password', password)
        page.click('button[type="submit"]')
        time.sleep(2)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)

            context = browser.new_context(viewport={'width': 1400, 'height': 900})
            page = context.new_page()
            print("Logging in as admin...")
            login(page, ADMIN_EMP, ADMIN_PASS)

            for route, fname, title, desc in PAGES:
                take(page, route, fname, title, desc)

            for route, fname, title, desc, sel, who in TABS:
                if who == 'admin':
                    print(f"  Capturing tab: {fname}...")
                    try:
                        page.goto(f'{BASE_URL}{route}', timeout=10000, wait_until='networkidle')
                        time.sleep(1)
                        page.click(sel)
                        time.sleep(1)
                        path = os.path.join(OUT_DIR, f'{fname}.png')
                        page.screenshot(path=path, full_page=True)
                        screenshots.append((path, fname, title, desc))
                    except Exception as e:
                        print(f"  WARN: Failed tab {fname}: {e}")
                        screenshots.append((None, fname, title, desc))

            context.close()

            context2 = browser.new_context(viewport={'width': 1400, 'height': 900})
            page2 = context2.new_page()
            print("Logging in as employee...")
            login(page2, 'EMP002', ADMIN_PASS)

            for route, fname, title, desc, sel, who in TABS:
                if who == 'employee':
                    print(f"  Capturing tab: {fname}...")
                    try:
                        page2.goto(f'{BASE_URL}{route}', timeout=10000, wait_until='networkidle')
                        time.sleep(1)
                        page2.click(sel)
                        time.sleep(1)
                        path = os.path.join(OUT_DIR, f'{fname}.png')
                        page2.screenshot(path=path, full_page=True)
                        screenshots.append((path, fname, title, desc))
                    except Exception as e:
                        print(f"  WARN: Failed tab {fname}: {e}")
                        screenshots.append((None, fname, title, desc))

            context2.close()
            browser.close()
    finally:
        proc.terminate()
        proc.wait()

    print(f"Captured {len([s for s in screenshots if s[0] is not None])} screenshots")
    return screenshots

DEPT_GROUPS = [
    ('Core & Authentication', 'Login, sessions, and access control', ['login']),
    ('Admin Dashboard', 'Real-time monitoring, module grid, agent productivity', ['admin_dashboard', 'admin_dashboard_productivity']),
    ('Employee Dashboard', 'Self-service home with module cards, breaks, activity', ['user_dashboard_modules', 'user_dashboard_breaks', 'user_dashboard_activity']),
    ('HR Administration', 'Employee master, org structure, holidays, import, audit', ['admin_users', 'org_chart', 'holidays', 'import_users', 'audit', 'reports']),
    ('Leave & Attendance', 'Leave applications, balance tracking, regularization', ['leaves', 'regularization']),
    ('Asset Management', 'Issue, track, and return company assets', ['assets']),
    ('Recruitment (ATS)', 'Jobs, candidates, hiring pipeline', ['jobs', 'candidates']),
    ('Onboarding & Offboarding', 'New-hire tasks and exit process', ['onboarding']),
    ('Payroll', 'Salary structures, payroll runs, deductions (PF/ESI/PT)', ['payroll', 'salary']),
    ('Performance Management', 'Goals, reviews, ratings, 360 feedback', ['goals', 'reviews', 'my_goals']),
    ('Expense Management', 'Claims, approvals, reimbursements', ['expenses_admin', 'expenses_emp']),
    ('Help Desk', 'Tickets, comments, assignments', ['tickets_admin', 'tickets_emp']),
    ('Document Management', 'Upload, download, manage employee documents', ['documents_admin', 'documents_emp']),
    ('Employee Self-Service', 'Profile, dependents, emergency contacts', ['profile']),
    ('Analytics & Reports', 'Charts, attrition risk, data export', ['analytics']),
]

DEPT_COLORS = [
    RGBColor(0x7C, 0x3A, 0xED), RGBColor(0x08, 0x91, 0xB2), RGBColor(0x05, 0x96, 0x69),
    RGBColor(0xDC, 0x26, 0x26), RGBColor(0x0D, 0x94, 0x88), RGBColor(0xE1, 0x1D, 0x48),
    RGBColor(0x02, 0x84, 0xC7), RGBColor(0x65, 0xA3, 0x0D), RGBColor(0x63, 0x66, 0xF1),
    RGBColor(0xD9, 0x77, 0x0C), RGBColor(0x84, 0x4B, 0xC7), RGBColor(0x14, 0x78, 0x6B),
    RGBColor(0xB9, 0x1D, 0x3B), RGBColor(0x43, 0x64, 0xEE), RGBColor(0xEA, 0x58, 0x0C),
]

def build_ppt(screenshots):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    W = Inches(13.333)
    H = Inches(7.5)

    DARK = RGBColor(0x1E, 0x29, 0x3B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN = RGBColor(0x19, 0x87, 0x54)
    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFB)
    GRAY = RGBColor(0x66, 0x66, 0x66)

    def add_background(slide, color=WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    def add_dept_heading(slide, text):
        add_shape(slide, Inches(0), Inches(0), W, Inches(1.1), DARK)
        add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                     text, font_size=36, bold=True, color=WHITE)

    def add_screenshot_slide(slide, img_path, title, desc):
        add_dept_heading(slide, title)
        add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                     desc, font_size=14, color=GRAY)
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.4))
            except Exception:
                add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                             f'[Screenshot]\n{title}', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                         f'[Screenshot not available]\n{title}', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

    shot_map = {s[1]: (s[0], s[2], s[3]) for s in screenshots}

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    add_shape(slide, Inches(0), H - Inches(0.15), W, Inches(0.15), GREEN)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 'HRMS Department Modules', font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 'Human Resource Management System', font_size=24, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                 'All Modules Organized by Department', font_size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 f'Generated: {datetime.now().strftime("%B %d, %Y")}', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_dept_heading(slide, 'Departments')
    lines = []
    for i, (dname, _, _) in enumerate(DEPT_GROUPS, 1):
        lines.append(f'  {i:02d}.  {dname}')
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                 '\n'.join(lines), font_size=20, color=DARK)

    for idx, (dept_name, dept_desc, fnames) in enumerate(DEPT_GROUPS):
        color = DEPT_COLORS[idx % len(DEPT_COLORS)]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, color)
        add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                     dept_name, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                     dept_desc, font_size=20, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

        for fname in fnames:
            if fname in shot_map:
                img_path, title, desc = shot_map[fname]
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_background(slide)
                add_screenshot_slide(slide, img_path, title, desc)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), GREEN)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 'Thank You', font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                 'HRMS Complete HR Solution', font_size=24, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 f'Generated {datetime.now().strftime("%Y-%m-%d %H:%M")}', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    prs.save(PPT_OUT)
    print(f"\nPowerPoint saved: {PPT_OUT}")

if __name__ == '__main__':
    screenshots = capture_all()
    build_ppt(screenshots)
    print("Done!")
